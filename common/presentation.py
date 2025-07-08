from pathlib import Path

from pptx import Presentation
from pptx.slide import SlideMaster, SlideLayout, Slide


class SinglePresentation:

    def __init__(self, ref_pptx: Path):
        self.__ref_pptx: Presentation = Presentation(ref_pptx.as_posix())
        self.__target_master: SlideMaster = self.__ref_pptx.slide_masters[0]

        self.__song_layout: SlideLayout = self.__target_master.slide_layouts[0]
        self.__none_layout: SlideLayout = self.__target_master.slide_layouts[1]
        self.delete_all()

    def add_song_layout(self, title: str, text: str):
        slide: Slide = self.__ref_pptx.slides.add_slide(slide_layout=self.__song_layout)
        for shape in slide.placeholders:
            print(f"=> {shape.placeholder_format.idx}|{shape.name}")

        print("==================================\n")
        slide.shapes.placeholders[0].text = title
        slide.placeholders[10].text = text
        # slide.placeholders[10].s
        # print(slide.placeholders[10].name)
        # slide.shapes.placeholders[10].text = title

    def add_none_layout(self):
        slide = self.__ref_pptx.slides.add_slide(slide_layout=self.__none_layout)

    def delete_all(self):

        slides = self.__ref_pptx.slides._sldIdLst

        # Create a list of rIds to drop relationships
        # This is important to truly remove the slide's content from the file.
        rIds_to_drop = []
        for slide_id_element in slides:
            rIds_to_drop.append(slide_id_element.rId)

        # Clear the slide list in the XML
        slides.clear()

        # Drop relationships to the deleted slides
        # This cleans up the internal file structure and prevents orphaned parts.
        for rId in rIds_to_drop:
            try:
                self.__ref_pptx.part.drop_rel(rId)
            except KeyError:
                # Handle cases where a relationship might already be gone or not found
                # (e.g., if a slide was already somewhat malformed or previously processed)
                pass

    def save(self, output_path: Path):
        self.__ref_pptx.save(output_path.as_posix())
