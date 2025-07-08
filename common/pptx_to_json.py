import json
from pptx import Presentation


def pptx_to_json(pptx_path, json_path):
    """
    Converts a PPTX file to a simplified JSON representation.
    It primarily extracts text content, slide layouts, and basic shape info.
    Font details are NOT explicitly stored in JSON, aiming to use layout defaults on conversion back.
    """
    prs = Presentation(pptx_path)
    presentation_data = {
        "slides": []
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }

        for shape_idx, shape in enumerate(slide.shapes):
            shape_data = {
                "shape_index": shape_idx,
                "shape_type": shape.shape_type.name,
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches,
            }

            if shape.has_text_frame:
                shape_data["has_text_frame"] = True
                shape_data["text"] = shape.text_frame.text
                # Note: We are not storing explicit font information here.
                # The assumption is that on re-creation, the layout will dictate the font.
            else:
                shape_data["has_text_frame"] = False

            # Add more specific shape handling if needed (e.g., pictures, tables)
            # This example focuses primarily on text shapes for simplicity.

            slide_data["shapes"].append(shape_data)
        presentation_data["slides"].append(slide_data)

    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, ensure_ascii=False, indent=4)
    print(f"PPTX converted to JSON successfully: {json_path}")




