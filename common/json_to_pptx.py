import json
from typing import List

from pptx import Presentation
from pptx.util import Inches

import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR


# (pptx_to_json function remains the same as it primarily extracts data)


def delete_all_slides(pptx_path):
    """
    Deletes all slides from a PowerPoint presentation.

    Args:
        pptx_path (str): The path to the PPTX file.
    """
    try:
        prs = Presentation(pptx_path)
    except FileNotFoundError:
        print(f"Error: File not found at {pptx_path}")
        return

    # Get the internal XML element for the slide list
    # _sldIdLst is a private attribute, but it's the most direct way
    # to manipulate the slide order/presence.
    slides = prs.slides._sldIdLst

    # Create a list of rIds to drop relationships
    # This is important to truly remove the slide's content from the file.
    rIds_to_drop = []
    for slide_id_element in slides:
        rIds_to_drop.append(slide_id_element.rId)

    # Clear the slide list in the XML
    slides.clear()

    # Drop relationships to the deleted slides
    # This cleans up the internal file structure and prevents orphaned parts.
    for rId in rIds_to_drop:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            # Handle cases where a relationship might already be gone or not found
            # (e.g., if a slide was already somewhat malformed or previously processed)
            pass

            # Save the modified presentation
    prs.save(pptx_path)
    print(f"All slides deleted from {pptx_path}")


def json_to_pptx(json_path, pptx_output_path):
    """
    Converts a JSON representation back to a PPTX file.
    It attempts to use the specified slide layouts and relies on them for font and formatting.
    This version tries to be more robust against 'Duplicate name' warnings.
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        presentation_data = json.load(f)

    # Create a new presentation. This automatically comes with a set of default slide masters and layouts.
    delete_all_slides(pptx_output_path)
    prs = Presentation()

    # Create a mapping of layout names to layout objects for quick lookup
    # We explicitly use the layouts from the *new* presentation.
    # all_layout:list = [*slide_master.slide_layout for slide_master in prs.slide_masters]
    layout_map = {layout.name: layout for layout in prs.slide_layouts}
    print(f"-->\n{json.dumps(layout_map,indent=4, default=str)}")

    for slide_data in presentation_data["slides"]:
        layout_name = slide_data["layout_name"]
        print(f"layout_namec: {layout_name}")

        assert layout_name in layout_map, f"layout({layout_name})이 존재하지 않습니다."
        slide_layout = None

        if layout_name in layout_map:
            slide_layout = layout_map[layout_name]
        else:
            print(f"Warning: Layout '{layout_name}' not found in the default new presentation. "
                  f"Using the 'Title and Content' layout as fallback (index 1 is common).")
            # Fallback to a common layout like 'Title and Content' if the exact named layout isn't present
            # or try to find a blank one.
            slide_layout = prs.slide_layouts[1]  # Common for Title and Content

        prs.slides.add_slide(slide_layout)

    prs.save(pptx_output_path)
    print(f"JSON converted to PPTX successfully: {pptx_output_path}")