import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

# (pptx_to_json function remains the same as it primarily extracts data)

def json_to_pptx(json_path, pptx_output_path):
    """
    Converts a JSON representation back to a PPTX file.
    It attempts to use the specified slide layouts and relies on them for font and formatting.
    This version tries to be more robust against 'Duplicate name' warnings.
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        presentation_data = json.load(f)

    # Create a new presentation. This automatically comes with a set of default slide masters and layouts.
    prs = Presentation()

    # Create a mapping of layout names to layout objects for quick lookup
    # We explicitly use the layouts from the *new* presentation.
    layout_map = {layout.name: layout for layout in prs.slide_layouts}

    for slide_data in presentation_data["slides"]:
        layout_name = slide_data["layout_name"]

        slide_layout = None
        if layout_name in layout_map:
            slide_layout = layout_map[layout_name]
        else:
            print(f"Warning: Layout '{layout_name}' not found in the default new presentation. "
                  f"Using the 'Title and Content' layout as fallback (index 1 is common).")
            # Fallback to a common layout like 'Title and Content' if the exact named layout isn't present
            # or try to find a blank one.
            try:
                slide_layout = prs.slide_layouts[1] # Common for Title and Content
            except IndexError:
                slide_layout = prs.slide_layouts[0] # Fallback to first layout

        slide = prs.slides.add_slide(slide_layout)

        for shape_data in slide_data["shapes"]:
            if shape_data["has_text_frame"]:
                text_content = shape_data["text"]

                # Try to find a suitable placeholder to put the text into.
                # This is crucial for inheriting layout-defined fonts and formatting.
                text_placed_in_placeholder = False
                for placeholder in slide.placeholders:
                    # A very basic heuristic: if it's a text placeholder and doesn't have text yet.
                    # In a more advanced scenario, you'd want to store placeholder index/type in JSON.
                    if placeholder.has_text_frame and not placeholder.text: # Check if it's empty
                        placeholder.text = text_content
                        text_placed_in_placeholder = True
                        break # Only use one placeholder per shape in JSON

                # If no suitable placeholder was found, add a generic text box.
                # Note: Generic text boxes might not inherit all layout styles as strongly as placeholders.
                if not text_placed_in_placeholder:
                    # Provide default dimensions if they are missing or for generic boxes
                    left = Inches(shape_data["left"]) if shape_data["left"] is not None else Inches(1)
                    top = Inches(shape_data["top"]) if shape_data["top"] is not None else Inches(1)
                    width = Inches(shape_data["width"]) if shape_data["width"] is not None else Inches(6)
                    height = Inches(shape_data["height"]) if shape_data["height"] is not None else Inches(1)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    txBox.text_frame.text = text_content
                    # You might want to explicitly set some basic font properties here if not using placeholders
                    # For example:
                    # from pptx.enum.text import MSO_AUTO_SIZE
                    # txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    # txBox.text_frame.word_wrap = True

    prs.save(pptx_output_path)
    print(f"JSON converted to PPTX successfully: {pptx_output_path}")

# --- Example Usage (same as before) ---
if __name__ == "__main__":
    print("Creating a dummy PPTX for testing...")
    prs_dummy = Presentation()

    # Add a title slide
    title_slide_layout = prs_dummy.slide_layouts[0] # Usually Title Slide
    slide1 = prs_dummy.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "My Sample Presentation"
    subtitle.text = "Created with Python"

    # Add a title and content slide
    bullet_slide_layout = prs_dummy.slide_layouts[1] # Usually Title and Content
    slide2 = prs_dummy.slides.add_slide(bullet_slide_layout)
    title2 = slide2.shapes.title
    body = slide2.placeholders[1]
    title2.text = "Features"
    tf = body.text_frame
    p = tf.add_paragraph()
    p.text = "This is a feature."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Sub-feature one."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sub-feature two."
    p.level = 1

    # Add a blank slide with a custom text box
    blank_slide_layout = prs_dummy.slide_layouts[6] # Usually Blank
    slide3 = prs_dummy.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1.0)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Custom text on blank slide."

    dummy_pptx_path = "sample_presentation.pptx"
    prs_dummy.save(dummy_pptx_path)
    print(f"Dummy PPTX created: {dummy_pptx_path}")

    # 2. Convert PPTX to JSON
    json_output_path = "output.json"
    pptx_to_json(dummy_pptx_path, json_output_path)

    # 3. Convert JSON back to PPTX
    reconverted_pptx_path = "reconverted_presentation.pptx"
    json_to_pptx(json_output_path, reconverted_pptx_path)

    print("\nProcess completed. Check 'output.json' and 'reconverted_presentation.pptx'.")
    print("Please be aware of the limitations regarding perfect layout/font preservation.")