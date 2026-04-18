import os
from typing import List
from pptx import Presentation
from src.core.schema import LayoutSchema, ShapeSchema, SlideSchema

class PPTAnalyzeError(Exception):
    pass

class PPTGenerateError(Exception):
    pass

def analyze_ppt(file_path: str) -> List[LayoutSchema]:
    if not os.path.exists(file_path):
        raise PPTAnalyzeError(f"File not found: {file_path}")
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise PPTAnalyzeError(f"Failed to parse PPT file: {str(e)}")

    layouts = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        shapes_schema = []
        for shape in layout.placeholders:
            if shape.has_text_frame:
                shapes_schema.append(ShapeSchema(
                    shape_name=shape.name,
                    shape_id=str(shape.shape_id),
                    shape_type="TEXT",
                    default_text=shape.text.replace('\n', ' ') if shape.text else ""
                ))
        
        if shapes_schema:
            layouts.append(LayoutSchema(
                layout_name=layout.name,
                layout_index=idx,
                shapes=shapes_schema
            ))
            
    return layouts

def generate_ppt(template_path: str, user_deck: List[SlideSchema], output_path: str) -> bool:
    if not os.path.exists(template_path):
        raise PPTGenerateError(f"Template file not found: {template_path}")
        
    try:
        prs = Presentation(template_path)
    except Exception as e:
        raise PPTGenerateError(f"Failed to parse template file: {str(e)}")

    for slide_data in user_deck:
        if slide_data.target_layout_index >= len(prs.slide_layouts):
            raise PPTGenerateError(f"Invalid layout index: {slide_data.target_layout_index}")
            
        slide_layout = prs.slide_layouts[slide_data.target_layout_index]
        new_slide = prs.slides.add_slide(slide_layout)
        
        for shape_data in slide_data.shapes:
            if shape_data.user_input is None:
                continue
                
            for shape in new_slide.placeholders:
                if shape.name == shape_data.shape_name and shape.has_text_frame:
                    shape.text = shape_data.user_input
                    break

    try:
        prs.save(output_path)
    except Exception as e:
        raise PPTGenerateError(f"Failed to save PPT: {str(e)}")
        
    return True
